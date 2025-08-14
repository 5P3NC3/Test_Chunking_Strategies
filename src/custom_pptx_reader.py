from pathlib import Path
from typing import Dict, List, Optional
from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document
from pptx import Presentation


class CustomPptxReader(BaseReader):
    """Custom PowerPoint reader using python-pptx."""
    
    def __init__(self):
        """Initialize the CustomPptxReader."""
        super().__init__()
    
    def load_data(
        self,
        file: Path,
        extra_info: Optional[Dict] = None,
    ) -> List[Document]:
        """Load data from PowerPoint file.
        
        Args:
            file: Path to the PowerPoint file
            extra_info: Additional information to add to the document
            
        Returns:
            List of documents
        """
        import pptx
        
        # Load the presentation
        prs = Presentation(str(file))
        
        # Extract text from all slides
        text_list = []
        
        for slide_number, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                    
                # Handle tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            if slide_text:
                # Add slide separator
                text_list.append(f"--- Slide {slide_number} ---")
                text_list.extend(slide_text)
        
        # Join all text
        full_text = "\n\n".join(text_list)
        
        # Create metadata
        metadata = extra_info or {}
        metadata.update({
            "file_name": file.name,
            "file_type": "pptx",
            "slide_count": len(prs.slides)
        })
        
        # Create and return document
        return [Document(text=full_text, metadata=metadata)]